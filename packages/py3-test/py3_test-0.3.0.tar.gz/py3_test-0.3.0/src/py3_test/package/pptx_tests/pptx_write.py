"""
@author: lijc210@163.com
@file: 1.py
@time: 2020/03/26
@desc: 功能描述。
"""
from pptx import Presentation
from pptx.util import Pt

prs = Presentation()  # 初始化一个ppt

# 设置PPT为16：9的宽屏
prs.slide_height = 6858000  # 设置ppt的高度
prs.slide_width = 12192000  # 设置ppt的宽度

# 用内置模板添加一个全空的ppt页面，内置貌似是7种，对应0-6
slide = prs.slides.add_slide(prs.slide_layouts[6])

# 全屏插入一张图片，图片最好提前处理长宽比，因为是不锁比例拉伸。
pic = slide.shapes.add_picture("123.png", 0, 0, width=prs.slide_width)

# 将刚插入的图片至于底层
slide.shapes._spTree.insert(1, pic._element)

# 在指定位置插入一个文本框，我按比例填的。
tBox = slide.shapes.add_textbox(
    left=prs.slide_width * 0.1,
    top=prs.slide_height * 0.1,
    width=prs.slide_width * 0.8,
    height=prs.slide_height * 0.8,
)
# 格式化为文本格式
tf = tBox.text_frame
# 运行插入
p = tf.add_paragraph()
# 设置粗体
p.font.bold = True
# 设置字体，一旦有中文就不正常，英文还好
# p.font.name = '楷体'
# 设置字体大小
p.font.size = Pt(40)
# 设置文本内容
p.text = "是什么呀 setthe z这是"
# 保存为文件
prs.save("test.pptx")
