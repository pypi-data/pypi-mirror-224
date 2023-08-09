"""
@author: lijc210@163.com
@file: pptx_read.py
@time: 2020/03/26
@desc: 功能描述。
"""
import json

from pptx import Presentation

shop_name_id_dict = {}
data_list = json.load(open("huizong.json", encoding="utf-8"))
for adict in data_list:
    shop_name_id_dict[adict["店铺名称"] + "-" + str(adict["店铺id"])] = adict["分站"]

data_dict = json.load(open("data.json", encoding="utf-8"))


def process(shop_name_id, city):
    print(shop_name_id)
    ppt = Presentation("商家体检报告-上海-金令设计装潢_tmp.pptx")
    for page, slide in enumerate(ppt.slides, 1):
        # print(slide)
        if page == 3:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    for row in table.table.rows:
                        if row.cells[0].text_frame.text == "锁单量":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["1锁单量"]
                            )
                        elif row.cells[0].text_frame.text == "分单量":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["2分单量"]
                            )
                        elif row.cells[0].text_frame.text == "锁单":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["3锁单"]
                            )
                        elif row.cells[0].text_frame.text == "不锁单":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["4不锁单"]
                            )
                        elif row.cells[0].text_frame.text == "赠单":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["5赠单"]
                            )
                        elif row.cells[0].text_frame.text == "补单":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["6补单"]
                            )
                        elif row.cells[0].text_frame.text == "接单完成率":
                            row.cells[1].text_frame.text = (
                                str(
                                    round(
                                        data_dict["锁单接单结构"][shop_name_id][0]["7接单完成率"]
                                        * 100,
                                        1,
                                    )
                                )
                                + "%"
                            )
                        elif row.cells[0].text_frame.text == "结转下月单量":
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构"][shop_name_id][0]["8结转下月单量"]
                            )
            # if shape.shape_type == 17:
            # 	text_frame = shape.text_frame
            # 	print(text_frame.text)
        elif page == 4:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    if data_dict["锁单接单结构-订单区域分布"].get(shop_name_id):
                        data_len = len(data_dict["锁单接单结构-订单区域分布"][shop_name_id])
                        print(shop_name_id, data_len)
                        for i, row in enumerate(table.table.rows):
                            if i == 0 or i > data_len:
                                continue
                            row.cells[0].text_frame.text = str(
                                data_dict["锁单接单结构-订单区域分布"][shop_name_id][i - 1]["订单区县"]
                            )
                            row.cells[1].text_frame.text = str(
                                data_dict["锁单接单结构-订单区域分布"][shop_name_id][i - 1]["订单量"]
                            )
                            row.cells[2].text_frame.text = (
                                str(
                                    round(
                                        data_dict["锁单接单结构-订单区域分布"][shop_name_id][i - 1][
                                            "该区县订单占比"
                                        ]
                                        * 100,
                                        1,
                                    )
                                )
                                + "%"
                            )
        elif page == 5:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    if data_dict["接单效率"].get(shop_name_id):
                        for i, row in enumerate(table.table.rows):
                            if i == 1:
                                row.cells[0].text_frame.text = str(
                                    data_dict["接单效率"][shop_name_id][0]["1订单量"]
                                )
                                row.cells[1].text_frame.text = str(
                                    data_dict["接单效率"][shop_name_id][0]["2接单量"]
                                )
                                if data_dict["接单效率"][shop_name_id][0]["3接单率"] == 1:
                                    row.cells[2].text_frame.text = "100%"
                                else:
                                    row.cells[2].text_frame.text = (
                                        str(
                                            round(
                                                data_dict["接单效率"][shop_name_id][0][
                                                    "3接单率"
                                                ]
                                                * 100,
                                                2,
                                            )
                                        )
                                        + "%"
                                    )

                                if (
                                    data_dict["接单效率"][shop_name_id][0]["4平均接单相应时长/min"]
                                    != "-"
                                ):
                                    if (
                                        data_dict["接单效率"][shop_name_id][0][
                                            "4平均接单相应时长/min"
                                        ]
                                        == 1
                                    ):
                                        row.cells[3].text_frame.text = str(100) + "分钟"
                                    else:
                                        row.cells[3].text_frame.text = (
                                            str(
                                                round(
                                                    data_dict["接单效率"][shop_name_id][0][
                                                        "4平均接单相应时长/min"
                                                    ],
                                                    2,
                                                )
                                            )
                                            + "分钟"
                                        )
                                else:
                                    row.cells[3].text_frame.text = "-"

                                if data_dict["接单效率"][shop_name_id][0]["6订单拨打率"] == 1:
                                    row.cells[4].text_frame.text = (
                                        str(
                                            data_dict["接单效率"][shop_name_id][0][
                                                "5虚拟号拨打订单量"
                                            ]
                                        )
                                        + "("
                                        + str(100)
                                        + "%"
                                        + ")"
                                    )
                                else:
                                    row.cells[4].text_frame.text = (
                                        str(
                                            data_dict["接单效率"][shop_name_id][0][
                                                "5虚拟号拨打订单量"
                                            ]
                                        )
                                        + "("
                                        + str(
                                            round(
                                                data_dict["接单效率"][shop_name_id][0][
                                                    "6订单拨打率"
                                                ]
                                                * 100,
                                                1,
                                            )
                                        )
                                        + "%"
                                        + ")"
                                    )
        elif page == 6:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    if data_dict["接单效率-接单员"].get(shop_name_id):
                        for _i1, row in enumerate(table.table.rows):
                            for j, dict in enumerate(
                                data_dict["接单效率-接单员"][shop_name_id][0:5]
                            ):
                                if row.cells[0].text_frame.text == "指标":
                                    row.cells[j + 1].text_frame.text = str(
                                        dict["1员工姓名"]
                                    )
                                elif row.cells[0].text_frame.text == "接单量":
                                    row.cells[j + 1].text_frame.text = str(dict["2接单量"])
                                elif row.cells[0].text_frame.text == "转派/接订单量":
                                    row.cells[j + 1].text_frame.text = str(
                                        dict["3转派/接订单量"]
                                    )
                                elif row.cells[0].text_frame.text == "实际需拨打订单量":
                                    row.cells[j + 1].text_frame.text = str(
                                        dict["4实际需拨打订单量"]
                                    )
                                elif "响应时长均值" in row.cells[0].text_frame.text:
                                    row.cells[j + 1].text_frame.text = str(
                                        round(dict["5平均响应时长"], 2)
                                    )
                                elif "虚拟号拨打量" in row.cells[0].text_frame.text:
                                    if dict["7拨打率"] == 1:
                                        row.cells[j + 1].text_frame.text = (
                                            str(dict["6拨打订单量"])
                                            + "("
                                            + str(100)
                                            + "%"
                                            + ")"
                                        )
                                    else:
                                        row.cells[j + 1].text_frame.text = (
                                            str(dict["6拨打订单量"])
                                            + "("
                                            + str(round(dict["7拨打率"] * 100, 1))
                                            + "%"
                                            + ")"
                                        )
        elif page == 18:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    for row in table.table.rows:
                        if row.cells[0].text_frame.text == "上传案例的设计师数量":
                            tmp1 = data_dict["设计师-店铺"].get(shop_name_id, [{}])
                            row.cells[1].text_frame.text = str(
                                tmp1[0].get("上传案例的设计师数量", 0)
                            )
                        elif row.cells[0].text_frame.text == "【全区域案例top20%】\n的设计师数":
                            tmp2 = data_dict["设计师-店铺"].get(shop_name_id, [{}])
                            if tmp2[0].get("占比", 0) == 1:
                                row.cells[1].text_frame.text = (
                                    str(tmp2[0].get("top20%设计师数", 0))
                                    + "("
                                    + "100%"
                                    + ")"
                                )
                            else:
                                row.cells[1].text_frame.text = (
                                    str(tmp2[0].get("top20%设计师数", 0))
                                    + "("
                                    + str(round(tmp2[0].get("占比1", 0) * 100, 2))
                                    + "%"
                                    + ")"
                                )
                        elif row.cells[0].text_frame.text == "被点评设计师数":
                            row.cells[1].text_frame.text = "请从右边取"
                        elif row.cells[0].text_frame.text == "设计师整体好评率":
                            row.cells[1].text_frame.text = "请从右边取"
                        elif row.cells[0].text_frame.text == "设计师整体差评率":
                            row.cells[1].text_frame.text = "请从右边取"
        elif page == 19:
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    if data_dict["设计师-人-案例数"].get(shop_name_id):
                        data_len = len(data_dict["设计师-人-案例数"][shop_name_id])
                        for i, row in enumerate(table.table.rows):
                            if i == 0 or i > data_len or i > 5:
                                continue
                            row.cells[0].text_frame.text = str(
                                data_dict["设计师-人-案例数"][shop_name_id][i - 1]["设计师名称"]
                            )
                            row.cells[1].text_frame.text = str(
                                data_dict["设计师-人-案例数"][shop_name_id][i - 1]["vr"]
                            )
                            row.cells[2].text_frame.text = str(
                                data_dict["设计师-人-案例数"][shop_name_id][i - 1]["平面"]
                            )
                            row.cells[3].text_frame.text = str(
                                data_dict["设计师-人-案例数"][shop_name_id][i - 1]["总案例数"]
                            )
                            row.cells[4].text_frame.text = str(
                                data_dict["设计师-人-案例数"][shop_name_id][i - 1]["区域排名"]
                            )
        elif page == 20:
            table_list = []
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    table_list.append(table)
            if data_dict["设计师-点评-设计师"].get(shop_name_id):
                data_len = len(data_dict["设计师-点评-设计师"][shop_name_id])
                for i, row in enumerate(table_list[0].table.rows):
                    if i == 0 or i > data_len or i > 5:
                        continue
                    row.cells[0].text_frame.text = str(
                        data_dict["设计师-点评-设计师"][shop_name_id][i - 1]["设计师名称"]
                    )
                    if data_dict["设计师-点评-设计师"][shop_name_id][i - 1]["好评率"] == 1:
                        row.cells[1].text_frame.text = "100%"
                    else:
                        row.cells[1].text_frame.text = (
                            str(
                                round(
                                    data_dict["设计师-点评-设计师"][shop_name_id][i - 1]["好评率"]
                                    * 100,
                                    1,
                                )
                            )
                            + "%"
                        )
                    row.cells[2].text_frame.text = str(
                        data_dict["设计师-点评-设计师"][shop_name_id][i - 1]["好评数"]
                    )
                    row.cells[3].text_frame.text = str(
                        data_dict["设计师-点评-设计师"][shop_name_id][i - 1]["区域排名"]
                    )

                data_len = len(data_dict["设计师-点评-设计师(差)"][shop_name_id])
                for i, row in enumerate(table_list[1].table.rows):
                    if i == 0 or i > data_len or i > 5:
                        continue
                    row.cells[0].text_frame.text = str(
                        data_dict["设计师-点评-设计师(差)"][shop_name_id][i - 1]["设计师名称"]
                    )
                    row.cells[1].text_frame.text = (
                        str(
                            round(
                                data_dict["设计师-点评-设计师(差)"][shop_name_id][i - 1]["差评率"]
                                * 100,
                                1,
                            )
                        )
                        + "%"
                    )
                    row.cells[2].text_frame.text = str(
                        data_dict["设计师-点评-设计师(差)"][shop_name_id][i - 1]["差评数"]
                    )
                    row.cells[3].text_frame.text = str(
                        data_dict["设计师-点评-设计师(差)"][shop_name_id][i - 1]["区域排名"]
                    )
        elif page == 21:
            table_list = []
            for shape in slide.shapes:
                # print(shape.shape_type)
                if shape.shape_type == 19:
                    table = shape
                    table_list.append(table)
            if data_dict["设计师-投诉率"].get(shop_name_id):
                data_len = len(data_dict["设计师-投诉率"][shop_name_id])
                for i, row in enumerate(table_list[1].table.rows):
                    if i == 0 or i > data_len or i > 5:
                        continue
                    row.cells[0].text_frame.text = str(
                        data_dict["设计师-投诉率"][shop_name_id][i - 1]["设计师名称"]
                    )
                    row.cells[1].text_frame.text = str(
                        data_dict["设计师-投诉率"][shop_name_id][i - 1]["投诉用户数"]
                    )
                    row.cells[2].text_frame.text = str(
                        data_dict["设计师-投诉率"][shop_name_id][i - 1]["签约用户数"]
                    )
                    row.cells[3].text_frame.text = (
                        str(
                            round(
                                data_dict["设计师-投诉率"][shop_name_id][i - 1]["投诉率"] * 100,
                                1,
                            )
                        )
                        + "%"
                    )

                data_len = len(data_dict["设计师-投诉类别"][shop_name_id])
                for i, row in enumerate(table_list[0].table.rows):
                    if i == 0 or i > data_len or i > 5:
                        continue
                    row.cells[0].text_frame.text = str(
                        data_dict["设计师-投诉类别"][shop_name_id][i - 1]["投诉类别"]
                    )
                    row.cells[1].text_frame.text = str(
                        data_dict["设计师-投诉类别"][shop_name_id][i - 1]["投诉用户数"]
                    )
    ppt.save("临时PPT数据/商家体检报告-{}-{}.pptx".format(city, shop_name_id.split("-")[0]))


if __name__ == "__main__":
    for shop_name_id, city in shop_name_id_dict.items():
        # if shop_name_id in ["天怡美装饰-215227980","天怡美装饰-215226137", "菡萏怡景装饰-215230579"]:
        #     process(shop_name_id, city)
        process(shop_name_id, city)
